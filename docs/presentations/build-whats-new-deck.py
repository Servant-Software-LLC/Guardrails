# -*- coding: utf-8 -*-
"""Build the 'Guardrails: What's New' deck WITH illustrations.
   Same layout/palette as build_deck.py; adds 5 visuals on slides 1,4,6,7,9."""
import io, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

# ---- palette ----
BLUE      = RGBColor(0x01, 0x49, 0x9B)
BLUE_DK   = RGBColor(0x01, 0x33, 0x6E)
ORANGE    = RGBColor(0xF3, 0x6F, 0x21)
ORANGE_DK = RGBColor(0xC4, 0x5C, 0x15)
DARK      = RGBColor(0x1F, 0x29, 0x37)
SLATE     = RGBColor(0x37, 0x41, 0x51)
MUTED     = RGBColor(0x6B, 0x72, 0x80)
MUTED2    = RGBColor(0x94, 0xA3, 0xB8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_TINT = RGBColor(0xEF, 0xF6, 0xFF)
BLUE_TINT2= RGBColor(0xE0, 0xEC, 0xFB)
GRAY_TINT = RGBColor(0xF8, 0xFA, 0xFC)
ORG_TINT  = RGBColor(0xFF, 0xF7, 0xED)
ORG_TINT2 = RGBColor(0xFE, 0xF0, 0xE6)
BORDER    = RGBColor(0xE5, 0xE7, 0xEB)
BORDER2   = RGBColor(0xD1, 0xD5, 0xDB)
CODEBG    = RGBColor(0x1E, 0x29, 0x3B)
CODEFG    = RGBColor(0xE5, 0xE7, 0xEB)
CODEMUT   = RGBColor(0x94, 0xA3, 0xB8)
CODEORG   = RGBColor(0xFB, 0xB9, 0x8C)
RED       = RGBColor(0xDC, 0x26, 0x26)
RED_TINT  = RGBColor(0xFE, 0xE2, 0xE2)
GREEN     = RGBColor(0x05, 0x96, 0x69)
ARIAL = "Arial"; MONO = "Courier New"
ASSETS = r"C:\Dev AI\Guardrails\docs\presentations\assets\whats-new"

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def _noline(sh): sh.line.fill.background()

def rect(s, l, t, w, h, fill=None, line=None, line_w=0.75, rounded=False, radius=0.08):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: _noline(shp)
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp

def settext(shape, paras, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
            wrap=True, ml=0.06, mr=0.06, mt=0.03, mb=0.03):
    tf = shape.text_frame if hasattr(shape, "text_frame") else shape
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Inches(ml); tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(mb)
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get("align", align)
        if "space_after" in pa: p.space_after = Pt(pa["space_after"])
        if "space_before" in pa: p.space_before = Pt(pa["space_before"])
        if "line" in pa: p.line_spacing = pa["line"]
        for (text, size, color, bold, font, italic) in pa["runs"]:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tf

def R(text, size, color, bold=False, font=ARIAL, italic=False):
    return (text, size, color, bold, font, italic)

def tbox(s, l, t, w, h, paras, **kw):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    settext(b, paras, **kw); return b

def header(s, title, step=None, label=None):
    rect(s, 0, 0, 10, 0.9, fill=BLUE)
    rect(s, 0, 0.9, 10, 0.05, fill=ORANGE)
    runs = []
    if step: runs.append(R(step + "   ", 17, RGBColor(0xBF,0xD3,0xF0), True))
    runs.append(R(title, 21, WHITE, True))
    tbox(s, 0.49, 0.12, 8.6, 0.66, [{"runs": runs}], anchor=MSO_ANCHOR.MIDDLE)
    if label:
        chip = rect(s, 7.7, 0.28, 1.85, 0.36, fill=BLUE_DK, rounded=True, radius=0.5)
        settext(chip, [{"align": PP_ALIGN.CENTER,
                        "runs": [R(label, 9.5, RGBColor(0xCF,0xDE,0xF5), True)]}],
                anchor=MSO_ANCHOR.MIDDLE)

def footer(s, idx):
    rect(s, 0, 5.5, 10, 0.125, fill=BLUE)
    tbox(s, 0.49, 5.46, 5, 0.2, [{"runs": [R("Guardrails — What's New", 8, MUTED2, False)]}])
    tbox(s, 8.0, 5.46, 1.5, 0.2, [{"align": PP_ALIGN.RIGHT, "runs": [R(f"{idx} / 10", 8, MUTED2, False)]}])

def notes(s, text): s.notes_slide.notes_text_frame.text = text

def card(s, l, t, w, h, title, lines, fill=GRAY_TINT, accent=ORANGE,
         border=BORDER, title_color=None, body_color=SLATE, title_size=14.5,
         body_size=12, bullet=True):
    rect(s, l, t, w, h, fill=fill, line=border, line_w=1, rounded=True, radius=0.06)
    title_color = title_color or BLUE
    tbox(s, l+0.18, t+0.13, w-0.36, 0.3, [{"runs": [R(title, title_size, title_color, True)]}])
    rect(s, l+0.18, t+0.5, min(w-0.36, 1.1), 0.035, fill=accent)
    paras = []
    for ln in lines:
        if isinstance(ln, tuple):
            paras.append({"space_after": 4, "line": 1.05, "runs": list(ln)})
        else:
            pre = ("•  " if bullet else "")
            paras.append({"space_after": 4, "line": 1.05,
                          "runs": [R(pre + ln, body_size, body_color, False)]})
    tbox(s, l+0.18, t+0.62, w-0.36, h-0.78, paras)

def codebox(s, l, t, w, h, lines, size=10.5, bg=CODEBG):
    rect(s, l, t, w, h, fill=bg, rounded=True, radius=0.05)
    paras = []
    for ln in lines:
        if isinstance(ln, list): paras.append({"space_after": 0, "line": 1.04, "runs": ln})
        else: paras.append({"space_after": 0, "line": 1.04, "runs": [R(ln, size, CODEFG, False, MONO)]})
    settext(s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)),
            paras, anchor=MSO_ANCHOR.MIDDLE, ml=0.18, mr=0.12)

def tagline_banner(s, l, t, w, h, sub=None):
    rect(s, l, t, w, h, fill=BLUE, rounded=True, radius=0.06)
    rect(s, l, t, 0.12, h, fill=ORANGE, rounded=False)
    paras = [{"align": PP_ALIGN.CENTER, "space_after": 3, "line": 1.05,
              "runs": [R("“A prompt may propose — only a deterministic gate may certify.”",
                         20, WHITE, True)]}]
    if sub:
        paras.append({"align": PP_ALIGN.CENTER,
                      "runs": [R(sub, 11, RGBColor(0xBF,0xD3,0xF0), False, ARIAL, True)]})
    settext(rect(s, l, t, w, h, fill=None), paras, anchor=MSO_ANCHOR.MIDDLE)

def style_table(tbl, header_fill=BLUE, header_fg=WHITE, body_size=11.5,
                head_size=12, col0_bold=True, zebra=(WHITE, GRAY_TINT)):
    for ci, cell in enumerate(tbl.rows[0].cells):
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
    for ri, row in enumerate(tbl.rows):
        for ci, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            if ri > 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = zebra[ri % 2]
            tf = cell.text_frame; tf.word_wrap = True
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.name = ARIAL
                    if ri == 0:
                        r.font.size = Pt(head_size); r.font.bold = True; r.font.color.rgb = header_fg
                    else:
                        r.font.size = Pt(body_size)
                        r.font.color.rgb = BLUE if (ci == 0 and col0_bold) else SLATE
                        r.font.bold = (ci == 0 and col0_bold)

def add_table(s, l, t, w, h, data, col_w=None):
    rows, cols = len(data), len(data[0])
    gtbl = s.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = gtbl.table; tbl.first_row = True; tbl.horz_banding = False
    for ri, rowdata in enumerate(data):
        for ci, val in enumerate(rowdata): tbl.cell(ri, ci).text = val
    if col_w:
        for ci, cw in enumerate(col_w): tbl.columns[ci].width = Inches(cw)
    return tbl

def arrow(s, l, t, w=0.34, h=0.34, color=BLUE, size=22):
    tbox(s, l, t, w, h, [{"align": PP_ALIGN.CENTER, "runs": [R("→", size, color, True)]}],
         anchor=MSO_ANCHOR.MIDDLE)

# ---- image helpers ----
def _to_transparent(im):
    im = im.convert("RGBA")
    out = []
    for (r, g, b, a) in im.getdata():
        mx = max(r, g, b); mn = min(r, g, b)
        if mn >= 236 and (mx - mn) <= 14: out.append((255, 255, 255, 0))
        else: out.append((r, g, b, a))
    im.putdata(out)
    bbox = im.getchannel("A").getbbox()
    return im.crop(bbox) if bbox else im

def load_img(name, transparent=True):
    im = Image.open(os.path.join(ASSETS, name))
    im = _to_transparent(im) if transparent else im.convert("RGBA")
    bio = io.BytesIO(); im.save(bio, format="PNG"); bio.seek(0)
    return bio, im.size

def load_quadrant(name, row, col):
    im = Image.open(os.path.join(ASSETS, name)).convert("RGBA")
    w, h = im.size
    q = im.crop((col*w//2, row*h//2, (col+1)*w//2, (row+1)*h//2))
    q = _to_transparent(q)
    bio = io.BytesIO(); q.save(bio, format="PNG"); bio.seek(0)
    return bio, q.size

def place_contain(s, bio, size, l, t, w, h, halign='c', valign='m'):
    iw, ih = size; ar = iw/ih; bar = w/h
    if ar >= bar: nw = w; nh = w/ar
    else: nh = h; nw = h*ar
    nl = l + (w-nw)/2 if halign=='c' else (l if halign=='l' else l+(w-nw))
    nt = t + (h-nh)/2 if valign=='m' else (t if valign=='t' else t+(h-nh))
    s.shapes.add_picture(bio, Inches(nl), Inches(nt), Inches(nw), Inches(nh))
    return (nl, nt, nw, nh)

def place_cover(s, bio, size, l, t, w, h):
    iw, ih = size; ar = iw/ih; bar = w/h
    if ar >= bar: nh = h; nw = h*ar
    else: nw = w; nh = w/ar
    nl = l + (w-nw)/2; nt = t + (h-nh)/2
    s.shapes.add_picture(bio, Inches(nl), Inches(nt), Inches(nw), Inches(nh))

def set_fill_alpha(shape, alpha_pct):
    """alpha_pct = opacity 0..100 (100 = fully opaque)."""
    spPr = shape._element.spPr
    srgb = spPr.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    srgb.append(srgb.makeelement(qn('a:alpha'), {'val': str(int(round(alpha_pct*1000)))}))

# =====================================================================
# SLIDE 1 — Title (full-bleed hero)
# =====================================================================
s = slide()
rect(s, 0, 0, 10, 5.625, fill=WHITE)
hero, hsz = load_img("01-title-hero.png", transparent=True)
place_cover(s, hero, hsz, 0, 0, 10, 5.625)
scrim = rect(s, 0, 4.18, 10, 1.32, fill=WHITE)        # lower-third legibility scrim
set_fill_alpha(scrim, 82)
chip = rect(s, 0.55, 0.40, 3.95, 0.38, fill=ORANGE, rounded=True, radius=0.5)
settext(chip, [{"align": PP_ALIGN.CENTER,
                "runs": [R("WHAT’S NEW  ·  SINCE THE LUNCH & LEARN", 11, WHITE, True)]}],
        anchor=MSO_ANCHOR.MIDDLE)
tbox(s, 0.55, 0.86, 7.3, 0.7, [{"runs": [R("Guardrails: What’s New", 38, BLUE, True)]}])
rect(s, 0.55, 4.30, 1.4, 0.04, fill=ORANGE)
tbox(s, 0.55, 4.42, 9.0, 0.45,
     [{"runs": [R("From one verified task to a whole team of them — running at once", 18, SLATE, True)]}])
tbox(s, 0.55, 4.92, 9.0, 0.65,
     [{"space_after": 2, "runs": [R("David Maltby", 14, DARK, True),
        R("    ·    A few weeks on from the Lunch & Learn  ·  June 2026", 12, MUTED, False)]},
      {"runs": [R("github.com/Servant-Software-LLC/Guardrails  ·  ServantSoftware.Guardrails  v1.0.0-preview.19",
                  10.5, MUTED, False, MONO)]}])
notes(s, "Last time I introduced Guardrails — the harness that makes a single AI task prove it worked "
         "before it counts as done. That was a few weeks ago. Since then the project has moved a long way, "
         "and the headline is a big one: Guardrails can now run many tasks IN PARALLEL, safely, and still keep "
         "every one of the guarantees we talked about. Today is a short 'what's new' — I'll recap in one "
         "slide, then spend most of our time on the new parallel engine.")
footer(s, 1)

# =====================================================================
# SLIDE 2 — Where We Left Off
# =====================================================================
s = slide()
header(s, "Where We Left Off", label="RECAP · 60 SEC")
stages = [("①", "Write a Plan", "Markdown"),
          ("②", "Break It Down", "/plan-breakdown"),
          ("③", "Review", "/guardrails-review"),
          ("④", "Execute", "guardrails run")]
x = 0.49; bw = 2.02; gap = 0.22
for i, (num, name, sub) in enumerate(stages):
    c = rect(s, x, 1.18, bw, 0.95, fill=BLUE_TINT, line=BORDER, line_w=1, rounded=True, radius=0.08)
    settext(c, [{"align": PP_ALIGN.CENTER, "space_after": 1,
                 "runs": [R(num+"  ", 14, ORANGE, True), R(name, 13, BLUE, True)]},
                {"align": PP_ALIGN.CENTER, "runs": [R(sub, 10.5, MUTED, False, MONO)]}],
            anchor=MSO_ANCHOR.MIDDLE)
    if i < 3: arrow(s, x+bw-0.02, 1.18, w=gap+0.06, h=0.95, color=BLUE, size=20)
    x += bw + gap
pc = rect(s, 0.49, 2.42, 9.03, 0.62, fill=GRAY_TINT, line=BORDER, line_w=1, rounded=True, radius=0.12)
settext(pc, [{"align": PP_ALIGN.CENTER,
              "runs": [R("The core promise:  ", 13, SLATE, False),
                       R("every task is verifiable, retryable, resumable.", 13, BLUE, True),
                       R("   Action + deterministic guardrails that prove it worked.", 12, MUTED, False)]}],
        anchor=MSO_ANCHOR.MIDDLE)
tagrow = rect(s, 0.49, 3.34, 9.03, 1.78, fill=BLUE, rounded=True, radius=0.06)
rect(s, 0.49, 3.34, 0.12, 1.78, fill=ORANGE)
settext(tagrow, [{"space_after": 6, "runs": [R("WHAT CHANGED SINCE JUNE 15", 12, RGBColor(0xBF,0xD3,0xF0), True)]},
                 {"line": 1.12,
                  "runs": [R("We took that promise — ", 19, WHITE, False),
                           R("prove every task worked", 19, WHITE, True, ARIAL, True),
                           R(" — and made it hold while ", 19, WHITE, False),
                           R("many tasks run at the same time.", 19, RGBColor(0xFF,0xC9,0xA3), True)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.3, mr=0.3)
notes(s, "Quick refresher for anyone who wasn't here, and a reset for everyone who was. A plan becomes a DAG "
         "of tasks; each task does something and then has to pass executable checks — the guardrails. "
         "Failures retry with targeted feedback, and an unfixable failure halts honestly for a human instead of "
         "pretending it's done. Hold that picture — everything new keeps it true. The whole challenge of the "
         "last few weeks was: how do you run tasks concurrently without weakening a single one of those guarantees?")
footer(s, 2)

# =====================================================================
# SLIDE 3 — The Headline: Parallel Execution
# =====================================================================
s = slide()
header(s, "The Headline: Parallel Execution")
card(s, 0.49, 1.12, 4.42, 2.75, "Before — Serial",
     [], fill=RED_TINT, accent=RED, title_color=RED, border=RGBColor(0xF3,0xC6,0xC6))
chain = rect(s, 0.72, 1.78, 3.96, 1.95, fill=None)
settext(chain, [{"align": PP_ALIGN.CENTER, "space_after": 5, "line": 1.0,
                 "runs": [R("01  →  02  →  03  →  04", 17, SLATE, True, MONO)]},
                {"align": PP_ALIGN.CENTER, "space_before": 6,
                 "runs": [R("one task at a time", 13, MUTED, False, ARIAL, True)]},
                {"align": PP_ALIGN.CENTER, "space_before": 10,
                 "runs": [R("independent work waits its turn\nfor no reason", 11.5, MUTED, False)]}],
        anchor=MSO_ANCHOR.MIDDLE)
card(s, 5.09, 1.12, 4.42, 2.75, "Now — Parallel",
     [], fill=BLUE_TINT, accent=BLUE, title_color=BLUE, border=BLUE_TINT2)
chip2 = rect(s, 7.55, 1.2, 1.78, 0.34, fill=BLUE, rounded=True, radius=0.5)
settext(chip2, [{"align": PP_ALIGN.CENTER, "runs": [R("maxParallelism: 3", 10, WHITE, True, MONO)]}],
        anchor=MSO_ANCHOR.MIDDLE)
for i, lbl in enumerate(["01", "02", "03"]):
    bx = rect(s, 5.42+i*1.18, 1.85, 1.0, 0.5, fill=WHITE, line=BLUE, line_w=1.25, rounded=True, radius=0.16)
    settext(bx, [{"align": PP_ALIGN.CENTER, "runs": [R(lbl, 15, BLUE, True, MONO)]}], anchor=MSO_ANCHOR.MIDDLE)
tbox(s, 5.27, 2.42, 4.06, 0.3, [{"align": PP_ALIGN.CENTER,
     "runs": [R("three tasks in flight at once", 12.5, BLUE, True, ARIAL, True)]}])
b04 = rect(s, 6.42, 2.85, 1.0, 0.46, fill=GRAY_TINT, line=BORDER2, line_w=1, rounded=True, radius=0.16)
settext(b04, [{"align": PP_ALIGN.CENTER, "runs": [R("04", 13, MUTED, True, MONO)]}], anchor=MSO_ANCHOR.MIDDLE)
tbox(s, 5.27, 3.33, 4.06, 0.3, [{"align": PP_ALIGN.CENTER,
     "runs": [R("04 waits on its deps, then runs", 11, MUTED, False)]}])
bc = rect(s, 0.49, 4.06, 9.03, 1.06, fill=ORG_TINT, line=RGBColor(0xF6,0xD8,0xBE), line_w=1, rounded=True, radius=0.1)
settext(bc, [{"space_after": 3, "line": 1.06,
              "runs": [R("The DAG already knew which tasks were independent. ", 13, SLATE, False),
                       R("Now the scheduler USES that", 13, ORANGE_DK, True),
                       R(" — filling the parallel slots with work that doesn’t depend on each other.", 13, SLATE, False)]},
             {"runs": [R("Turn it on: ", 12, SLATE, False),
                       R("maxParallelism: 3", 12, BLUE, True, MONO),
                       R("  in guardrails.json.  Serial (", 12, SLATE, False),
                       R("1", 12, BLUE, True, MONO),
                       R(") stays the safe default — you opt in.", 12, SLATE, False)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.22, mr=0.22)
notes(s, "This is the big one. The dependency graph always told us which tasks could run side by side — the "
         "doc-update task never needed to wait for the feature task. Before, we ran them one at a time anyway. "
         "Now the scheduler runs up to three at once by default. Three is deliberate — impressive in a demo, "
         "gentle on disk. And notice: serial is still the default. You opt INTO parallel; you don't get surprised "
         "by it. The hard part wasn't running things at once — it was doing it WITHOUT two agents stepping on "
         "each other. That's the next three slides.")
footer(s, 3)

# =====================================================================
# SLIDE 4 — Worktrees Isolate (illustration left, text right)
# =====================================================================
s = slide()
header(s, "Worktrees Isolate", step="①")
tbox(s, 0.49, 0.98, 9.03, 0.36,
     [{"runs": [R("The danger with concurrency: two agents writing the same workspace → corruption. ", 12.5, SLATE, False),
                R("The fix is physical, not clever.", 12.5, ORANGE_DK, True)]}])
img, isz = load_img("05-worktree-isolation.png", transparent=True)
place_contain(s, img, isz, 0.5, 1.5, 3.95, 3.5)
tbox(s, 4.78, 1.55, 4.72, 1.0,
     [{"line": 1.12, "runs": [R("Each running task gets its own ", 12.5, SLATE, False),
        R("git worktree", 12.5, BLUE, True), R(" — a separate working copy on its own branch. "
        "Each writes ", 12.5, SLATE, False), R("only its own tree", 12.5, BLUE, True),
        R("; the harness merges results back one at a time.", 12.5, SLATE, False)]}])
k1 = rect(s, 4.78, 2.72, 4.74, 1.0, fill=BLUE, rounded=True, radius=0.08)
rect(s, 4.78, 2.72, 0.1, 1.0, fill=ORANGE)
settext(k1, [{"line": 1.08, "runs": [R("“Agent B clobbered agent A’s edits” is ", 12.5, WHITE, False),
             R("structurally impossible", 12.5, RGBColor(0xFF,0xC9,0xA3), True),
             R(" — not mitigated. Impossible.", 12.5, WHITE, False)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.22, mr=0.18)
k2 = rect(s, 4.78, 3.84, 4.74, 0.95, fill=GRAY_TINT, line=BORDER, line_w=1, rounded=True, radius=0.08)
settext(k2, [{"line": 1.08, "runs": [R("A linear chain ", 12.5, SLATE, False),
             R("reuses one worktree", 12.5, BLUE, True),
             R(", passed along — so we don’t pay for trees we don’t need.", 12.5, SLATE, False)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.2, mr=0.18)
notes(s, "The instinct is to solve concurrency with locks and cleverness. We didn't. Each task runs in its own "
         "physical copy of the repo — a git worktree — so there is simply no shared mutable workspace "
         "for two agents to fight over. The whole class of 'one agent overwrote the other's work' bug doesn't get "
         "reduced, it gets removed by construction. And we're careful about cost: if tasks form a straight line, "
         "they hand the same worktree down the chain instead of making a fresh copy each time. Git shares the "
         "history under the hood, so what multiplies is just the working files, not the whole repo.")
footer(s, 4)

# =====================================================================
# SLIDE 5 — The Write-Scope Check
# =====================================================================
s = slide()
header(s, "The Write-Scope Check", step="②")
tbox(s, 0.49, 1.0, 9.03, 0.62,
     [{"line": 1.08, "runs": [R("Each task can declare what it’s allowed to write — its ", 12.5, SLATE, False),
                R("writeScope", 12.5, BLUE, True, MONO),
                R(".  After the action runs, a ", 12.5, SLATE, False),
                R("deterministic, read-only check", 12.5, ORANGE_DK, True),
                R(" (just a ", 12.5, SLATE, False), R("git diff", 12.5, BLUE, True, MONO),
                R(" membership test) confirms the task only touched files inside that surface.", 12.5, SLATE, False)]}])
tbl = add_table(s, 0.49, 1.74, 9.03, 2.05,
    [["It buys us", "How"],
     ["Small, clean diffs", "a task that strays outside its scope fails early and locally — a retry, not a tangled conflict later"],
     ["TDD enforced for free", "the implementation task’s scope EXCLUDES the test files — it literally cannot edit the tests to pass"],
     ["Easier merges later", "narrow, non-overlapping diffs rarely conflict when two parallel branches finally meet"]],
    col_w=[2.5, 6.53])
style_table(tbl, body_size=11.5, head_size=13)
note = rect(s, 0.49, 3.95, 9.03, 1.17, fill=ORG_TINT, line=RGBColor(0xF6,0xD8,0xBE), line_w=1, rounded=True, radius=0.08)
settext(note, [{"space_after": 4, "line": 1.06,
                "runs": [R("It writes nothing and reverts nothing across tasks — it’s a pure check.", 12, ORANGE_DK, True)]},
               {"line": 1.06,
                "runs": [R("On a violation, only the out-of-scope files are rolled back, and the agent retries with ", 12, SLATE, False),
                         R("“you touched files you weren’t allowed to — here’s the list.”", 12, SLATE, True, ARIAL, True)]},
               {"space_before": 4, "line": 1.04,
                "runs": [R("This replaced the older hash-capture machinery for protecting test files — same guarantee, far simpler.", 10.5, MUTED, False)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.22, mr=0.22)
notes(s, "This replaced an older, fiddlier mechanism for protecting test files with something much simpler: just "
         "declare what a task is allowed to write, then deterministically check the diff stayed inside it. The "
         "elegant part is what it gives us almost for free. The TDD rule we care about — the code can't quietly "
         "edit the tests to make them pass — becomes a one-line scope: the implementation task's write-scope "
         "excludes the test files, so editing them is a guardrail failure, full stop. And small, well-scoped diffs "
         "are the thing that makes the next slide — merging concurrent work — tractable.")
footer(s, 5)

# =====================================================================
# SLIDE 6 — A Prompt Proposes, a Gate Certifies (gate emblem)
# =====================================================================
s = slide()
header(s, "A Prompt Proposes, a Gate Certifies", step="③")
flow = [("git auto-merge", "tries first", BLUE, BLUE_TINT, BLUE),
        ("AI resolves conflict", "BYTE producer — never the verdict", ORANGE, ORG_TINT, ORANGE),
        ("Deterministic re-verify", "build + tests on merged bytes = the verdict", BLUE, BLUE_TINT, BLUE)]
fx = 0.49; fw = 2.84; fgap = 0.26
for i, (t1, t2, ac, fl, ln) in enumerate(flow):
    fb = rect(s, fx, 1.12, fw, 1.16, fill=fl, line=ln, line_w=1.25, rounded=True, radius=0.08)
    settext(fb, [{"align": PP_ALIGN.CENTER, "space_after": 3,
                  "runs": [R(t1, 13.5, ac if ac==ORANGE else BLUE, True)]},
                 {"align": PP_ALIGN.CENTER, "line": 1.04, "runs": [R(t2, 11, SLATE, False)]}],
            anchor=MSO_ANCHOR.MIDDLE)
    if i < 2: arrow(s, fx+fw-0.04, 1.12, w=fgap+0.08, h=1.16, color=MUTED, size=20)
    fx += fw + fgap
sub = rect(s, 0.49, 2.42, 9.03, 0.72, fill=GRAY_TINT, line=BORDER, line_w=1, rounded=True, radius=0.1)
settext(sub, [{"line": 1.06, "align": PP_ALIGN.CENTER,
               "runs": [R("Two deterministic checks gate the AI’s bytes — ", 11.5, SLATE, False),
                        R("no conflict markers left; only the conflicted files touched", 11.5, BLUE, True),
                        R(".  The AI never gets a softer gate. Can’t pass → ", 11.5, SLATE, False),
                        R("honest halt for a human.", 11.5, ORANGE_DK, True)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.2, mr=0.2)
emb, esz = load_img("02-gate-emblem.png", transparent=True)
place_contain(s, emb, esz, 0.5, 3.22, 1.75, 1.32)
tagline_banner(s, 2.4, 3.28, 7.12, 1.18,
               sub="A final whole-repo gate runs once at the end and backstops every merge.")
notes(s, "This is the heart of the whole design, and the one sentence I'd want you to walk out with. We let an AI "
         "do something genuinely useful — resolve a merge conflict, which is fiddly, judgement-heavy work. But "
         "we never let the AI decide that it worked. Its output is just bytes. Those bytes only count after "
         "deterministic checks confirm it didn't cheat — no leftover conflict markers, no touching files "
         "outside the conflict — and then the real verdict is the same build-and-test re-verify that a clean "
         "automatic merge would face. The AI gets no discount. If it can't pass, the run halts honestly for a "
         "human. That principle — a prompt may propose, only a deterministic gate may certify — is what "
         "lets us use AI aggressively AND trust the result. It's the same idea as the guardrails themselves, now "
         "applied to merging.")
footer(s, 6)

# =====================================================================
# SLIDE 7 — The Breakdown Skill Got Smarter (4 icons)
# =====================================================================
s = slide()
header(s, "The Breakdown Skill Got Smarter")
tbox(s, 0.49, 0.98, 9.03, 0.3,
     [{"runs": [R("/plan-breakdown", 12.5, BLUE, True, MONO),
                R(" now generates better task folders out of the box:", 12.5, SLATE, False)]}])
grid = [("Tests-first by default", ORANGE,
         "Writes the author-the-tests task BEFORE the implementation task. TDD is the default doctrine, not an afterthought."),
        ("No silent guesses", BLUE,
         "If the test framework is ambiguous it ASKS you — instead of quietly picking one and surprising you."),
        ("Servers & executables get wired up", ORANGE,
         "Inserts an entry-point-wiring task and a live smoke test: “does it actually start and respond?” — not just unit tests."),
        ("UI plans get real UI tasks", BLUE,
         "Implementation tasks plus UI-presence guardrails, so a “build the screen” plan can’t pass with an empty screen.")]
cw, ch = 4.42, 1.62; gx, gy = 0.49, 1.34; gxp, gyp = 0.19, 0.16
for i, (title, ac, body) in enumerate(grid):
    col = i % 2; row = i // 2
    l = gx + col*(cw+gxp); t = gy + row*(ch+gyp)
    fill = ORG_TINT if ac==ORANGE else BLUE_TINT
    brd = RGBColor(0xF6,0xD8,0xBE) if ac==ORANGE else BLUE_TINT2
    tcol = ORANGE_DK if ac==ORANGE else BLUE
    rect(s, l, t, cw, ch, fill=fill, line=brd, line_w=1, rounded=True, radius=0.06)
    ico, isz2 = load_quadrant("03-feature-icons.png", row, col)
    place_contain(s, ico, isz2, l+0.15, t+0.16, 0.66, 0.66)
    tbox(s, l+0.95, t+0.16, cw-1.12, 0.6,
         [{"line": 1.0, "runs": [R(title, 13.5, tcol, True)]}], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, l+0.95, t+0.78, 1.0, 0.03, fill=ac)
    tbox(s, l+0.18, t+0.88, cw-0.36, ch-0.98,
         [{"line": 1.05, "runs": [R(body, 11.5, SLATE, False)]}])
notes(s, "The other half of the work went into the skill that turns your plan into the task folder. It's gotten a "
         "lot more opinionated in good ways. It defaults to test-driven — tests get authored first, as their "
         "own task. It refuses to guess silently: if it can't tell which test framework you use, it asks rather "
         "than picking one and surprising you. For anything that runs as a server or an executable, it now adds a "
         "task that actually wires up the entry point and a smoke test that boots it and checks it responds — "
         "because 'it compiles' isn't 'it runs'. And for UI work it generates real UI tasks with checks that the "
         "UI is actually present. The theme: push more correctness into the generated structure, before a human "
         "even reviews it.")
footer(s, 7)

# =====================================================================
# SLIDE 8 — Sharper Feedback & Clearer Eyes
# =====================================================================
s = slide()
header(s, "Sharper Feedback & Clearer Eyes")
tbox(s, 0.49, 0.98, 9.03, 0.3,
     [{"runs": [R("The everyday-experience upgrades that make all of the above usable:", 12.5, SLATE, False)]}])
tbl = add_table(s, 0.49, 1.36, 9.03, 3.8,
    [["Upgrade", "Why it matters"],
     ["Verify, don’t replay", "guardrails check the action’s recorded output — they don’t blindly re-run it and hope for the same answer"],
     ["Streaming logs", "the agent’s transcript and output stream to disk live — watch a task think in real time"],
     ["Interactive DAG diagram", "diagram.html — pan, zoom, full-screen, and click through to each task"],
     ["Clickable log links", "the run prints links straight to each task’s post-mortem logs"],
     ["Honest budget halt", "a per-run cost cap (maxCostUsd) stops the run cleanly when the dollar budget is hit"]],
    col_w=[2.55, 6.48])
style_table(tbl, body_size=11.5, head_size=13)
notes(s, "None of this is glamorous, but it's what makes the engine pleasant to live with. Guardrails now verify "
         "the output the action actually produced rather than re-running things and hoping for determinism. Logs "
         "stream to disk as the agent works, so you can watch a long task unfold instead of staring at a spinner. "
         "The dependency graph is now an interactive HTML diagram you can pan, zoom, and click into. The run hands "
         "you clickable links straight to the logs for any task that needs a look. And there's a real cost cap — "
         "set a dollar budget, and when it's reached the run stops honestly instead of quietly burning money. Small "
         "things, but they're the difference between a demo and a tool you'd actually leave running.")
footer(s, 8)

# =====================================================================
# SLIDE 9 — The One Idea (balance image)
# =====================================================================
s = slide()
header(s, "The One Idea")
tagline_banner(s, 0.49, 1.04, 9.03, 1.0)
tbl = add_table(s, 0.49, 2.18, 6.15, 2.4,
    [["New capability", "The proposer  (AI)", "The certifier  (deterministic)"],
     ["Parallel execution", "agents work concurrently", "worktrees isolate; the harness merges"],
     ["AI merge-resolution", "AI resolves the conflict", "2 checks + build/test re-verify on merged bytes"],
     ["Write-scope", "the action edits files", "a git diff membership check"],
     ["Smarter breakdown", "the skill proposes tasks", "human review + the guardrails still gate"]],
    col_w=[1.5, 1.9, 2.75])
style_table(tbl, body_size=10, head_size=11)
bal, bsz = load_img("04-balance.png", transparent=True)
place_contain(s, bal, bsz, 6.78, 2.12, 2.74, 2.5)
tagline2 = rect(s, 0.49, 4.66, 9.03, 0.6, fill=ORANGE, rounded=True, radius=0.12)
settext(tagline2, [{"align": PP_ALIGN.CENTER,
                    "runs": [R("We let AI do more — by trusting it less.  ", 17, WHITE, True),
                             R("That’s the whole trick.", 15, WHITE, False, ARIAL, True)]}],
        anchor=MSO_ANCHOR.MIDDLE)
notes(s, "If there's a single thread through all of this, it's that line. We gave AI agents MORE freedom this "
         "cycle — to run in parallel, to resolve merge conflicts, to scaffold richer task structures. And the "
         "way we made that safe was not to trust the AI more — it was to put a cheap, deterministic gate in "
         "front of every place the AI proposes something. More autonomy on one side, more verification on the "
         "other. That balance is the entire philosophy of the project, and the parallel engine is just the biggest "
         "example of it so far.")
footer(s, 9)

# =====================================================================
# SLIDE 10 — Getting Started & Questions
# =====================================================================
s = slide()
header(s, "Getting Started & Questions")
tbox(s, 0.49, 1.04, 4.5, 0.3, [{"runs": [R("Update to the latest preview", 12.5, BLUE, True)]}])
codebox(s, 0.49, 1.36, 4.5, 1.0,
        [[R("dotnet tool update --global \\", 9.5, CODEFG, False, MONO)],
         [R("    ServantSoftware.Guardrails ", 9.5, CODEFG, False, MONO),
          R("--prerelease", 9.5, CODEORG, False, MONO)],
         [R("guardrails skills install", 9.5, CODEFG, False, MONO)],
         [R("# refresh /plan-breakdown & /guardrails-review", 9, CODEMUT, False, MONO)]])
tbox(s, 0.49, 2.5, 4.5, 0.3, [{"runs": [R("Go parallel — in guardrails.json", 12.5, BLUE, True)]}])
codebox(s, 0.49, 2.82, 4.5, 0.62,
        [[R("{ ", 11, CODEFG, False, MONO), R("\"maxParallelism\"", 11, CODEORG, False, MONO),
          R(": ", 11, CODEFG, False, MONO), R("3", 11, RGBColor(0x7E,0xE7,0xC7), False, MONO),
          R(" }", 11, CODEFG, False, MONO)]])
tbox(s, 0.49, 3.54, 4.5, 0.95,
     [{"space_after": 3, "line": 1.05,
       "runs": [R("Requires git (worktree mode). Serial runs (", 10.5, MUTED, False),
                R("1", 10.5, SLATE, True, MONO),
                R(") need no git and keep the old model.", 10.5, MUTED, False)]},
      {"space_before": 4, "runs": [R("ServantSoftware.Guardrails  ·  v1.0.0-preview.19", 10, SLATE, True, MONO)]},
      {"runs": [R("github.com/Servant-Software-LLC/Guardrails", 10, MUTED, False, MONO)]}])
rect(s, 5.25, 1.04, 4.27, 4.06, fill=BLUE_TINT, line=BLUE_TINT2, line_w=1, rounded=True, radius=0.05)
tbox(s, 5.45, 1.2, 3.9, 0.36, [{"runs": [R("Takeaways", 15, BLUE, True)]}])
rect(s, 5.45, 1.6, 1.1, 0.035, fill=ORANGE)
takeaways = [
    ("Parallel now", "up to 3 tasks at once, isolated by git worktrees"),
    ("AI merges conflicts", "but only behind deterministic gates"),
    ("Smarter breakdown", "TDD-first, smoke tests, UI tasks, no silent guesses"),
    ("One principle throughout", "a prompt may propose, only a deterministic gate may certify")]
yy = 1.78
for i, (h, b) in enumerate(takeaways):
    nb = rect(s, 5.45, yy+0.02, 0.34, 0.34, fill=ORANGE, rounded=True, radius=0.5)
    settext(nb, [{"align": PP_ALIGN.CENTER, "runs": [R(str(i+1), 12, WHITE, True)]}], anchor=MSO_ANCHOR.MIDDLE)
    tbox(s, 5.92, yy-0.04, 3.5, 0.78,
         [{"space_after": 1, "runs": [R(h, 12.5, DARK, True)]},
          {"line": 1.02, "runs": [R(b, 11, SLATE, False)]}])
    yy += 0.82
notes(s, "To try it: update the tool, re-install the skills so you get the smarter breakdown, and add one line to "
         "your config to go parallel. The only new requirement is git, and only when you actually turn parallelism "
         "on — serial runs are unchanged. I'm dogfooding all of this on real work right now, including using "
         "Guardrails to build Guardrails. Happy to take questions, or to sit down and break down one of your real "
         "plans afterward.")
footer(s, 10)

out = r"C:\Dev AI\Guardrails\docs\presentations\guardrails-whats-new-2026-06-22.pptx"
prs.save(out)
print("Saved:", out)
