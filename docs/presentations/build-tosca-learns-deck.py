# -*- coding: utf-8 -*-
"""Build the 'Are we missing a step in Agentic Engineering?' deck (Tosca Learns, 2026-06-30).
   Same palette/house-style as build-whats-new-deck.py; typographic (no image assets)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette (identical to the What's New deck) ----
BLUE      = RGBColor(0x01, 0x49, 0x9B)
BLUE_DK   = RGBColor(0x01, 0x33, 0x6E)
ORANGE    = RGBColor(0xF3, 0x6F, 0x21)
ORANGE_DK = RGBColor(0xC4, 0x5C, 0x15)
DARK      = RGBColor(0x1F, 0x29, 0x37)
SLATE     = RGBColor(0x37, 0x41, 0x51)
MUTED     = RGBColor(0x6B, 0x72, 0x80)
MUTED2    = RGBColor(0x94, 0xA3, 0xB8)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_TINT = RGBColor(0xEF, 0xF6, 0xFF)
BLUE_TINT2= RGBColor(0xE0, 0xEC, 0xFB)
LTBLUE    = RGBColor(0xDC, 0xE7, 0xF6)
GRAY_TINT = RGBColor(0xF8, 0xFA, 0xFC)
ORG_TINT  = RGBColor(0xFF, 0xF7, 0xED)
BORDER    = RGBColor(0xE5, 0xE7, 0xEB)
BORDER2   = RGBColor(0xD1, 0xD5, 0xDB)
CODEBG    = RGBColor(0x1E, 0x29, 0x3B)
CODEFG    = RGBColor(0xE5, 0xE7, 0xEB)
CODEMUT   = RGBColor(0x94, 0xA3, 0xB8)
CODEORG   = RGBColor(0xFB, 0xB9, 0x8C)
CODEGRN   = RGBColor(0x7E, 0xE7, 0xC7)
RED       = RGBColor(0xDC, 0x26, 0x26)
RED_TINT  = RGBColor(0xFE, 0xE2, 0xE2)
RED_BRD   = RGBColor(0xF3, 0xC6, 0xC6)
GREEN     = RGBColor(0x05, 0x96, 0x69)
ARIAL = "Arial"; MONO = "Courier New"

prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def _noline(sh): sh.line.fill.background()

def rect(s, l, t, w, h, fill=None, line=None, line_w=0.75, rounded=False, radius=0.08):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: _noline(shp)
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp

def settext(shape, paras, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT,
            wrap=True, ml=0.06, mr=0.06, mt=0.03, mb=0.03):
    tf = shape.text_frame if hasattr(shape, "text_frame") else shape
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Inches(ml); tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(mb)
    for i, pa in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = pa.get("align", align)
        if "space_after" in pa: p.space_after = Pt(pa["space_after"])
        if "space_before" in pa: p.space_before = Pt(pa["space_before"])
        if "line" in pa: p.line_spacing = pa["line"]
        for (text, size, color, bold, font, italic) in pa["runs"]:
            r = p.add_run(); r.text = text
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic; r.font.name = font
    return tf

def R(text, size, color, bold=False, font=ARIAL, italic=False):
    return (text, size, color, bold, font, italic)

def tbox(s, l, t, w, h, paras, **kw):
    b = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    settext(b, paras, **kw); return b

def header(s, title, step=None, label=None):
    rect(s, 0, 0, 10, 0.9, fill=BLUE)
    rect(s, 0, 0.9, 10, 0.05, fill=ORANGE)
    runs = []
    if step: runs.append(R(step + "   ", 17, RGBColor(0xBF,0xD3,0xF0), True))
    runs.append(R(title, 21, WHITE, True))
    tbox(s, 0.49, 0.12, 7.5, 0.66, [{"runs": runs}], anchor=MSO_ANCHOR.MIDDLE)
    if label:
        chip = rect(s, 7.7, 0.28, 1.85, 0.36, fill=BLUE_DK, rounded=True, radius=0.5)
        settext(chip, [{"align": PP_ALIGN.CENTER,
                        "runs": [R(label, 10, RGBColor(0xCF,0xDE,0xF5), True)]}],
                anchor=MSO_ANCHOR.MIDDLE)

def footer(s, idx):
    rect(s, 0, 5.5, 10, 0.125, fill=BLUE)
    tbox(s, 0.49, 5.46, 6, 0.2, [{"runs": [R("Tosca Learns  ·  Agentic Engineering", 8, MUTED2, False)]}])
    tbox(s, 8.0, 5.46, 1.5, 0.2, [{"align": PP_ALIGN.RIGHT, "runs": [R(f"{idx} / 12", 8, MUTED2, False)]}])

def notes(s, text): s.notes_slide.notes_text_frame.text = text

def codebox(s, l, t, w, h, lines, size=10.5, bg=CODEBG):
    rect(s, l, t, w, h, fill=bg, rounded=True, radius=0.05)
    paras = []
    for ln in lines:
        if isinstance(ln, list): paras.append({"space_after": 0, "line": 1.18, "runs": ln})
        else: paras.append({"space_after": 0, "line": 1.18, "runs": [R(ln, size, CODEFG, False, MONO)]})
    settext(s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)),
            paras, anchor=MSO_ANCHOR.MIDDLE, ml=0.2, mr=0.12)

def tagline_banner(s, l, t, w, h, sub=None):
    rect(s, l, t, w, h, fill=BLUE, rounded=True, radius=0.06)
    rect(s, l, t, 0.12, h, fill=ORANGE, rounded=False)
    paras = [{"align": PP_ALIGN.CENTER, "space_after": 3, "line": 1.05,
              "runs": [R("“A prompt may propose — only a deterministic gate may certify.”",
                         20, WHITE, True)]}]
    if sub:
        paras.append({"align": PP_ALIGN.CENTER,
                      "runs": [R(sub, 11, RGBColor(0xBF,0xD3,0xF0), False, ARIAL, True)]})
    settext(rect(s, l, t, w, h, fill=None), paras, anchor=MSO_ANCHOR.MIDDLE)

def style_table(tbl, header_fill=BLUE, header_fg=WHITE, body_size=11.5,
                head_size=12, col0_bold=True, zebra=(WHITE, GRAY_TINT)):
    for ci, cell in enumerate(tbl.rows[0].cells):
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
    for ri, row in enumerate(tbl.rows):
        for ci, cell in enumerate(row.cells):
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05); cell.margin_bottom = Inches(0.05)
            if ri > 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = zebra[ri % 2]
            tf = cell.text_frame; tf.word_wrap = True
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.name = ARIAL
                    if ri == 0:
                        r.font.size = Pt(head_size); r.font.bold = True; r.font.color.rgb = header_fg
                    else:
                        r.font.size = Pt(body_size)
                        r.font.color.rgb = BLUE if (ci == 0 and col0_bold) else SLATE
                        r.font.bold = (ci == 0 and col0_bold)

def add_table(s, l, t, w, h, data, col_w=None):
    rows, cols = len(data), len(data[0])
    gtbl = s.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
    tbl = gtbl.table; tbl.first_row = True; tbl.horz_banding = False
    for ri, rowdata in enumerate(data):
        for ci, val in enumerate(rowdata): tbl.cell(ri, ci).text = val
    if col_w:
        for ci, cw in enumerate(col_w): tbl.columns[ci].width = Inches(cw)
    return tbl

def arrow(s, l, t, w=0.34, h=0.34, color=BLUE, size=22, ch="→"):
    tbox(s, l, t, w, h, [{"align": PP_ALIGN.CENTER, "runs": [R(ch, size, color, True)]}],
         anchor=MSO_ANCHOR.MIDDLE)

def fieldnote(s, l, t, w, h, runs_list):
    """orange-tinted 'from real runs' callout"""
    rect(s, l, t, w, h, fill=ORG_TINT, line=RGBColor(0xF6,0xD8,0xBE), line_w=1, rounded=True, radius=0.08)
    settext(rect(s, l, t, w, h, fill=None), [{"line": 1.06, "runs": runs_list}],
            anchor=MSO_ANCHOR.MIDDLE, ml=0.22, mr=0.2)

# =====================================================================
# SLIDE 1 — Title
# =====================================================================
s = slide()
rect(s, 0, 0, 10, 5.625, fill=WHITE)
rect(s, 0, 0, 10, 0.16, fill=BLUE)
rect(s, 0, 0.16, 10, 0.04, fill=ORANGE)
chip = rect(s, 0.62, 1.05, 5.1, 0.42, fill=ORANGE, rounded=True, radius=0.5)
settext(chip, [{"align": PP_ALIGN.CENTER,
                "runs": [R("AGENTIC ENGINEERING   ·   TOSCA LEARNS", 12, WHITE, True)]}],
        anchor=MSO_ANCHOR.MIDDLE)
tbox(s, 0.58, 1.68, 9.0, 1.9,
     [{"line": 1.04, "runs": [R("Are we missing a step", 43, BLUE, True)]},
      {"line": 1.04, "runs": [R("in Agentic Engineering?", 43, BLUE, True)]}])
rect(s, 0.64, 3.55, 1.5, 0.05, fill=ORANGE)
tbox(s, 0.6, 3.72, 9.0, 0.7,
     [{"line": 1.12, "runs": [R("An honest question about how we ship AI-written code — ", 16, SLATE, True),
                              R("not a pitch.", 16, ORANGE_DK, True)]}])
tbox(s, 0.6, 4.62, 9.0, 0.6,
     [{"runs": [R("David Maltby", 14, DARK, True),
                R("     ·     Tosca Learns   ·   June 2026", 12, MUTED, False)]}])
notes(s, "I want to ask a question, not sell anything. We've all gotten very good at getting an AI agent "
         "to write code for us. I think there might be a step in the middle we've quietly skipped — and I "
         "want to spend ten minutes on what that step is and why it matters. At the very end I'll show you one "
         "thing I built that fills it, but honestly the question is the point, not my tool.")
footer(s, 1)

# =====================================================================
# SLIDE 2 — How we build with agents today
# =====================================================================
s = slide()
header(s, "How We Build With Agents Today", label="THE STATUS QUO")
boxes = [("write a plan", "in Markdown", "plan.md", BLUE_TINT, BLUE_TINT2, BLUE),
         ("hand the WHOLE file", "to an agent: “implement”", "one prompt · one shot", ORG_TINT, RGBColor(0xF6,0xD8,0xBE), ORANGE_DK),
         ("a big PR", "of changes", "merge & move on", BLUE_TINT, BLUE_TINT2, BLUE)]
bx = 0.49; bw = 2.86; bgap = 0.27
for i, (t1, t2, t3, fl, brd, ac) in enumerate(boxes):
    fb = rect(s, bx, 1.35, bw, 1.4, fill=fl, line=brd, line_w=1.25, rounded=True, radius=0.08)
    settext(fb, [{"align": PP_ALIGN.CENTER, "space_after": 2, "runs": [R(t1, 14.5, ac, True)]},
                 {"align": PP_ALIGN.CENTER, "space_after": 6, "line": 1.03, "runs": [R(t2, 11.5, SLATE, False)]},
                 {"align": PP_ALIGN.CENTER, "runs": [R(t3, 10, MUTED, False, MONO)]}],
            anchor=MSO_ANCHOR.MIDDLE)
    if i < 2: arrow(s, bx+bw-0.05, 1.35, w=bgap+0.1, h=1.4, color=MUTED, size=24)
    bx += bw + bgap
tbox(s, 0.49, 2.95, 9.03, 0.4, [{"align": PP_ALIGN.CENTER,
     "runs": [R("It’s fast. It often works. It feels like magic.", 14, SLATE, True, ARIAL, True)]}])
bc = rect(s, 0.49, 3.5, 9.03, 1.5, fill=BLUE, rounded=True, radius=0.06)
rect(s, 0.49, 3.5, 0.12, 1.5, fill=ORANGE)
settext(bc, [{"space_after": 5, "line": 1.1,
              "runs": [R("But step back: what did we actually ", 15, WHITE, False),
                       R("verify", 15, RGBColor(0xFF,0xC9,0xA3), True),
                       R("?", 15, WHITE, False)]},
             {"line": 1.1,
              "runs": [R("We reviewed the prose going in, and the diff coming out. Everything ", 13.5, LTBLUE, False),
                       R("in between", 13.5, WHITE, True, ARIAL, True),
                       R(" — the agent just ", 13.5, LTBLUE, False),
                       R("did.", 13.5, WHITE, True)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.3, mr=0.3)
notes(s, "Here's the workflow. You write a nice plan in markdown — maybe a really good one. You paste the "
         "whole thing into an agent and say 'implement this.' One prompt, one shot, and out comes a pull request. "
         "And it's genuinely impressive how often this works. But notice what we reviewed: the plan on the way in, "
         "and the diff on the way out. Everything in the middle — how the work was split up, in what order, what "
         "got checked — the agent did on its own, and we never saw it. The next three slides are about what lives "
         "in that gap.")
footer(s, 2)

# =====================================================================
# SLIDE 3 — Problem 1: the breakdown is invisible
# =====================================================================
s = slide()
header(s, "The Breakdown Is Invisible", step="①", label="PROBLEM 1 / 3")
tbox(s, 0.49, 1.05, 9.03, 0.42,
     [{"runs": [R("A markdown plan is ", 14, SLATE, False), R("prose", 14, ORANGE_DK, True),
                R(".  It doesn’t show a ", 14, SLATE, False), R("DAG", 14, BLUE, True),
                R(".", 14, SLATE, False)]}])
card = rect(s, 0.49, 1.6, 4.55, 2.35, fill=GRAY_TINT, line=BORDER, line_w=1, rounded=True, radius=0.06)
tbox(s, 0.69, 1.74, 4.2, 0.34, [{"runs": [R("So the agent silently decides:", 13, BLUE, True)]}])
rect(s, 0.69, 2.12, 1.1, 0.035, fill=ORANGE)
tbox(s, 0.69, 2.26, 4.2, 1.6,
     [{"space_after": 7, "line": 1.05, "runs": [R("•  how to ", 12.5, SLATE, False), R("split", 12.5, DARK, True), R(" the work into tasks", 12.5, SLATE, False)]},
      {"space_after": 7, "line": 1.05, "runs": [R("•  what ", 12.5, SLATE, False), R("order", 12.5, DARK, True), R(" to do them in", 12.5, SLATE, False)]},
      {"line": 1.05, "runs": [R("•  what ", 12.5, SLATE, False), R("depends", 12.5, DARK, True), R(" on what", 12.5, SLATE, False)]}])
k = rect(s, 5.22, 1.6, 4.3, 2.35, fill=BLUE, rounded=True, radius=0.06)
rect(s, 5.22, 1.6, 0.11, 2.35, fill=ORANGE)
settext(k, [{"space_after": 6, "line": 1.12,
             "runs": [R("That decomposition ", 14.5, WHITE, False), R("is the real engineering", 14.5, RGBColor(0xFF,0xC9,0xA3), True),
                      R(" —", 14.5, WHITE, False)]},
            {"line": 1.12, "runs": [R("and it’s ", 14.5, WHITE, False), R("hidden from review.", 14.5, WHITE, True),
                                    R("  You can’t review a plan-of-work that was never written down.", 12.5, LTBLUE, False)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.26, mr=0.22)
fieldnote(s, 0.49, 4.12, 9.03, 0.86,
          [R("From real runs:  ", 11.5, ORANGE_DK, True),
           R("a task compiled against a type a ", 11.5, SLATE, False),
           R("different", 11.5, DARK, True),
           R(" task hadn’t produced yet — a dependency nothing in the plan ever surfaced.", 11.5, SLATE, False)])
notes(s, "First problem. A plan in markdown is prose — it reads top to bottom. But real work is a graph: this "
         "task depends on that one, these two are independent, that one has to come last. When you one-shot a "
         "markdown file, the agent invents that graph in its head and never shows you. And the breakdown is "
         "arguably the most important engineering decision in the whole job — get the dependencies wrong and "
         "everything downstream is shaky. We literally hit this: one task was writing code that referenced a type "
         "another task was supposed to create first, and nothing in the plan made that ordering visible. You can't "
         "review a structure that doesn't exist on paper. "
         "(For your reference, not the audience's: that was issue #176 in our own tracker — a failure mode the "
         "gate caught by dogfooding. Numbers stay off the slides.)")
footer(s, 3)

# =====================================================================
# SLIDE 4 — Problem 2: the agent grades its own homework
# =====================================================================
s = slide()
header(s, "The Agent Grades Its Own Homework", step="②", label="PROBLEM 2 / 3")
tbox(s, 0.49, 1.05, 9.03, 0.7,
     [{"line": 1.1, "runs": [R("Who decides a task is ", 14, SLATE, False), R("done", 14, DARK, True),
                R("?  The agent.  “Definition of done” becomes whatever the model’s output ", 14, SLATE, False),
                R("claims", 14, ORANGE_DK, True),
                R(" — with no deterministic check that the claim is ", 14, SLATE, False),
                R("true.", 14, BLUE, True)]}])
gt = rect(s, 0.49, 1.9, 9.03, 0.5, fill=RED_TINT, line=RED_BRD, line_w=1, rounded=True, radius=0.1)
settext(gt, [{"align": PP_ALIGN.CENTER, "runs": [R("Green  ≠  correct", 16, RED, True)]}], anchor=MSO_ANCHOR.MIDDLE)
card1 = rect(s, 0.49, 2.56, 4.46, 1.42, fill=GRAY_TINT, line=BORDER, line_w=1, rounded=True, radius=0.06)
settext(card1, [{"space_after": 5, "line": 1.08,
                 "runs": [R("all the tests pass…", 14, DARK, True)]},
                {"line": 1.12, "runs": [R("over a feature that was ", 13, SLATE, False),
                                        R("never actually wired up.", 13, RED, True)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.2, mr=0.18)
card2 = rect(s, 5.06, 2.56, 4.46, 1.42, fill=GRAY_TINT, line=BORDER, line_w=1, rounded=True, radius=0.06)
settext(card2, [{"space_after": 5, "line": 1.08,
                 "runs": [R("a “tests pass” check…", 14, DARK, True)]},
                {"line": 1.12, "runs": [R("that passes on code that ", 13, SLATE, False),
                                        R("doesn’t even compile.", 13, RED, True)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.2, mr=0.18)
bn = rect(s, 0.49, 4.16, 9.03, 0.82, fill=BLUE, rounded=True, radius=0.08)
rect(s, 0.49, 4.16, 0.12, 0.82, fill=ORANGE)
settext(bn, [{"align": PP_ALIGN.CENTER, "line": 1.06,
              "runs": [R("A prompt is very good at telling you ", 14, WHITE, False),
                       R("what you want to hear.", 14, RGBColor(0xFF,0xC9,0xA3), True)]}],
        anchor=MSO_ANCHOR.MIDDLE)
notes(s, "Second problem, and it's the big one: determinism. In a one-shot, the thing that decides whether the "
         "work succeeded is the same model that did the work. Definition of done is whatever it says in its "
         "summary — 'Done! All tests pass.' But 'the agent says it's done' is not a proof. We've seen all-green "
         "test suites sitting on top of a feature that was never connected to anything — the tests passed because "
         "they tested the wrong layer. We've seen a 'tests pass' check report success on code that didn't compile. "
         "A language model is, by design, a machine for producing plausible-sounding output. If you let it certify "
         "its own work, plausible is all you're guaranteed. "
         "(Your reference: those were issues #120 and #155 in our tracker.)")
footer(s, 4)

# =====================================================================
# SLIDE 5 — Problem 3: it'll move the goalposts to pass
# =====================================================================
s = slide()
header(s, "…And It’ll Move the Goalposts to Pass", step="③", label="PROBLEM 3 / 3")
tbox(s, 0.49, 1.05, 9.03, 0.74,
     [{"space_after": 4, "line": 1.08,
       "runs": [R("Tell an agent ", 14, SLATE, False), R("“make the tests pass,”", 14, DARK, True, ARIAL, True),
                R(" and one of the easiest paths is to ", 14, SLATE, False), R("edit the test.", 14, RED, True)]},
      {"line": 1.08, "runs": [R("Or it quietly writes files ", 14, SLATE, False),
                R("outside the task it was given", 14, RED, True),
                R(" — touching things it shouldn’t.", 14, SLATE, False)]}])
warn = rect(s, 0.49, 2.5, 9.03, 0.62, fill=RED_TINT, line=RED_BRD, line_w=1, rounded=True, radius=0.1)
settext(warn, [{"align": PP_ALIGN.CENTER,
                "runs": [R("In a one-shot, ", 14, SLATE, False), R("nothing stops either.", 14, RED, True),
                         R("  You get a green check and a quiet lie.", 14, SLATE, False)]}],
        anchor=MSO_ANCHOR.MIDDLE)
fix = rect(s, 0.49, 3.36, 9.03, 1.62, fill=BLUE, rounded=True, radius=0.06)
rect(s, 0.49, 3.36, 0.12, 1.62, fill=ORANGE)
settext(fix, [{"space_after": 5, "runs": [R("THE FIX IS DETERMINISTIC, NOT HOPEFUL", 11.5, RGBColor(0xBF,0xD3,0xF0), True)]},
              {"line": 1.12, "runs": [R("The implementation task is ", 15, WHITE, False),
                R("only allowed to write the source", 15, RGBColor(0xFF,0xC9,0xA3), True),
                R(" — its scope ", 15, WHITE, False),
                R("excludes the test files", 15, WHITE, True),
                R(".  Editing a test to pass is a hard, mechanical failure.  ", 14, LTBLUE, False),
                R("No trust required.", 14, WHITE, True)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.3, mr=0.26)
notes(s, "Third problem follows straight from the second. If the agent both does the work and decides the work is "
         "done, then the path of least resistance can be to game the check. Told 'make the failing test pass,' "
         "editing the test is a perfectly valid way to make it pass — and a terrible way to make the code correct. "
         "Same with scope: an agent fixing one thing wanders off and edits five other files, and in a one-shot "
         "you'd never notice until something breaks later. Here's the thing — this isn't solved by a better prompt "
         "or a sterner instruction. It's solved by a deterministic rule: the implementation task is mechanically "
         "only allowed to touch the source files, not the tests. Edit a test and it's an automatic failure, full "
         "stop. No trust required. (Your reference: that scope-enforcement work was issue #136.)")
footer(s, 5)

# =====================================================================
# SLIDE 6 — Synthesis: a prompt proposes, nothing certifies
# =====================================================================
s = slide()
header(s, "A Prompt Proposes — Nothing Certifies")
tbox(s, 0.49, 1.0, 9.03, 0.32,
     [{"runs": [R("Line the three up and the same gap appears every time:", 13.5, SLATE, False)]}])
tbl = add_table(s, 0.49, 1.42, 9.03, 1.95,
    [["What we one-shot", "What’s missing"],
     ["The agent DECIDES the breakdown", "a visible DAG you can review"],
     ["The agent WRITES each task’s prompt in its head", "the prompts, out in the open"],
     ["The agent DECLARES it’s done", "a deterministic gate that proves it"]],
    col_w=[4.75, 4.28])
style_table(tbl, body_size=12.5, head_size=13)
band = rect(s, 0.49, 3.62, 9.03, 0.66, fill=ORG_TINT, line=RGBColor(0xF6,0xD8,0xBE), line_w=1, rounded=True, radius=0.1)
settext(band, [{"align": PP_ALIGN.CENTER,
                "runs": [R("We’ve automated the ", 14, SLATE, False), R("proposing.", 14, BLUE, True),
                         R("   We ", 14, SLATE, False), R("skipped the certifying.", 14, ORANGE_DK, True)]}],
        anchor=MSO_ANCHOR.MIDDLE)
q = rect(s, 0.49, 4.46, 9.03, 0.62, fill=BLUE, rounded=True, radius=0.08)
rect(s, 0.49, 4.46, 0.12, 0.62, fill=ORANGE)
settext(q, [{"align": PP_ALIGN.CENTER,
             "runs": [R("That skipped step is the question.", 17, WHITE, True)]}], anchor=MSO_ANCHOR.MIDDLE)
notes(s, "So here's the synthesis. Across all three problems it's the same shape: the agent proposes — it proposes "
         "a breakdown, it proposes the work, it proposes that it's finished — and nothing independent and "
         "deterministic ever certifies any of it. We've gotten incredibly good at the proposing half and we've "
         "quietly skipped the certifying half. That missing certifying step — a breakdown you can see, prompts you "
         "can read, and a gate that mechanically proves 'done' — that's the step I think we might be missing. "
         "Everything from here is just: what would filling it look like?")
footer(s, 6)

# =====================================================================
# SLIDE 7 — My first guess: Beads
# =====================================================================
s = slide()
header(s, "My First Guess: Beads", label="MARCH 2026")
qb = rect(s, 0.49, 1.18, 9.03, 0.78, fill=GRAY_TINT, line=BORDER, line_w=1, rounded=True, radius=0.08)
rect(s, 0.49, 1.18, 0.1, 0.78, fill=ORANGE)
settext(qb, [{"line": 1.06, "runs": [R("Beads", 15, BLUE, True),
             R("  —  “a distributed graph issue tracker for AI agents.”", 15, SLATE, False, ARIAL, True)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.25)
tbox(s, 0.49, 2.2, 9.03, 1.3,
     [{"space_after": 6, "line": 1.16,
       "runs": [R("A real ", 14, SLATE, False), R("dependency graph", 14, BLUE, True),
                R(" of tasks, persistent across long agent sessions, with ", 14, SLATE, False),
                R("ready", 14, DARK, True, MONO), R("-task detection.", 14, SLATE, False)]},
      {"line": 1.16, "runs": [R("Genuinely clever. ", 14, SLATE, False),
                R("I was hooked.", 14, ORANGE_DK, True)]}])
aside = rect(s, 0.49, 3.74, 9.03, 0.74, fill=ORG_TINT, line=RGBColor(0xF6,0xD8,0xBE), line_w=1, rounded=True, radius=0.1)
settext(aside, [{"align": PP_ALIGN.CENTER,
                 "runs": [R("Full disclosure: a certain ", 12.5, SLATE, False, ARIAL, True),
                          R("Gastown", 12.5, ORANGE_DK, True, ARIAL, True),
                          R(" obsession may have helped.", 12.5, SLATE, False, ARIAL, True)]}],
        anchor=MSO_ANCHOR.MIDDLE)
notes(s, "My first attempt at filling the gap wasn't to build anything — it was to adopt something. Back in March "
         "I found Beads, which bills itself as a distributed graph issue tracker for AI agents. And it's a genuinely "
         "smart project: instead of a flat markdown to-do list, your tasks live in a real dependency graph, it "
         "persists across long agent sessions, it can tell an agent which tasks are 'ready' because their blockers "
         "are done. I was sold. I'll admit the name didn't hurt — anyone who knows me knows I have a bit of a "
         "Gastown thing. So for a while I thought: this is the missing step. It wasn't — and the way it fell short "
         "is what told me what the step actually is.")
footer(s, 7)

# =====================================================================
# SLIDE 8 — Why Beads fell short
# =====================================================================
s = slide()
header(s, "Why Beads Fell Short (For Me)")
tbox(s, 0.49, 0.98, 9.03, 0.32,
     [{"runs": [R("It organized the work — but left the ", 13, SLATE, False),
                R("engineering judgment", 13, ORANGE_DK, True),
                R(" with the agent:", 13, SLATE, False)]}])
tbl = add_table(s, 0.49, 1.38, 9.03, 2.55,
    [["", "Beads"],
     ["1.  Who runs the tasks", "the AGENTS managed their own tasks"],
     ["2.  Can I see it", "tasks AND dependency chains were hard to visualize"],
     ["3.  Can I see the prompts", "no — the prompt used for each task wasn’t visible"],
     ["4.  Who defines “done”", "left to the AGENT performing the task"]],
    col_w=[3.4, 5.63])
style_table(tbl, body_size=12, head_size=12.5)
band = rect(s, 0.49, 4.12, 9.03, 0.86, fill=BLUE, rounded=True, radius=0.06)
rect(s, 0.49, 4.12, 0.12, 0.86, fill=ORANGE)
settext(band, [{"line": 1.1, "runs": [R("Same gap as the one-shot — just better organized.  ", 14, WHITE, False),
               R("The human still wasn’t in the loop where it counts.", 14, RGBColor(0xFF,0xC9,0xA3), True)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.3, mr=0.26)
notes(s, "Here's where it came up short, for my purposes. One: the agents managed the tasks — claimed them, "
         "updated them, closed them — so I was still trusting the agent to drive. Two: I couldn't easily see it. "
         "The tasks and especially the dependency chains were hard to visualize; there was no picture I could put "
         "in front of a reviewer. Three: I couldn't see the prompt that drove each task — the actual instruction "
         "was invisible. And four, the big one: definition of done was still left to the agent doing the task. So "
         "Beads organized the work beautifully, but it left every piece of engineering judgment with the agent. "
         "It's the same gap as the one-shot, just tidier. To be fair, Beads is solving a different problem — agent "
         "memory — it just wasn't the verification step I was missing. That's when I started building the step itself.")
footer(s, 8)

# =====================================================================
# SLIDE 9 — One answer: Guardrails
# =====================================================================
s = slide()
header(s, "One Answer: Guardrails", label="A TASTE")
tbox(s, 0.49, 0.98, 9.03, 0.32,
     [{"runs": [R("Not a prescription — just what I built to fill the gap.  The same four steps, but the ", 12.5, SLATE, False),
                R("middle is now visible and verified:", 12.5, BLUE, True)]}])
stages = [("①", "Write a Plan", "Markdown"),
          ("②", "Break It Down", "a VISIBLE DAG"),
          ("③", "Review", "a human signs off"),
          ("④", "Execute", "each task proves it")]
x = 0.49; bw = 2.02; gap = 0.22
for i, (num, name, sub) in enumerate(stages):
    c = rect(s, x, 1.4, bw, 0.92, fill=BLUE_TINT, line=BORDER, line_w=1, rounded=True, radius=0.08)
    settext(c, [{"align": PP_ALIGN.CENTER, "space_after": 1,
                 "runs": [R(num+"  ", 13, ORANGE, True), R(name, 12.5, BLUE, True)]},
                {"align": PP_ALIGN.CENTER, "runs": [R(sub, 9.5, MUTED, False)]}],
            anchor=MSO_ANCHOR.MIDDLE)
    if i < 3: arrow(s, x+bw-0.02, 1.4, w=gap+0.06, h=0.92, color=BLUE, size=18)
    x += bw + gap
tbl = add_table(s, 0.49, 2.56, 9.03, 2.42,
    [["The gap", "How this fills it"],
     ["breakdown was invisible", "a DAG of task folders you open and review"],
     ["prompts were hidden", "every task’s prompt is a file, in the open"],
     ["“done” was the agent’s word", "deterministic guardrails — executable pass/fail checks"],
     ["the agent owned the structure", "YOU review and sign off before anything runs"]],
    col_w=[3.5, 5.53])
style_table(tbl, body_size=11.5, head_size=12)
notes(s, "So this is the thing I built, and I'm going to keep it deliberately high-level because the mechanism "
         "matters less than the shape. It's the same four steps you'd expect — write a plan, break it down, review, "
         "execute — but the middle two, the part the one-shot hides, are now out in the open. The breakdown is a "
         "real DAG you can look at: a folder per task, with its dependencies. Each task's prompt is a plain file you "
         "can read and edit. And critically, 'done' for a task isn't the agent's opinion — it's a set of "
         "deterministic guardrails, ordinary executable checks that pass or fail. And before a single task runs, a "
         "human reviews the whole structure. Every row in that table is just one of the problems from earlier, "
         "turned around.")
footer(s, 9)

# =====================================================================
# SLIDE 10 — The one idea
# =====================================================================
s = slide()
header(s, "The One Idea")
tbox(s, 0.49, 1.12, 9.03, 0.34, [{"align": PP_ALIGN.CENTER,
     "runs": [R("If you remember one sentence:", 13.5, SLATE, False, ARIAL, True)]}])
tagline_banner(s, 0.7, 1.66, 8.6, 1.5)
mid = rect(s, 0.49, 3.5, 9.03, 1.5, fill=GRAY_TINT, line=BORDER, line_w=1, rounded=True, radius=0.08)
settext(mid, [{"align": PP_ALIGN.CENTER, "space_after": 8, "line": 1.1,
               "runs": [R("Let the AI do ", 18, DARK, False), R("more", 18, ORANGE_DK, True),
                        R(" — by ", 18, DARK, False), R("trusting it less.", 18, BLUE, True)]},
              {"align": PP_ALIGN.CENTER, "line": 1.08,
               "runs": [R("The agent proposes the work; a cheap, deterministic check decides whether it counts.", 13, SLATE, False)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.3, mr=0.3)
notes(s, "If you take one thing away, it's this line. A prompt may propose — only a deterministic gate may "
         "certify. The agent is allowed to be creative, fast, even brilliant — it proposes. But the thing that "
         "decides whether the work actually counts is never the prompt; it's a cheap, boring, deterministic check. "
         "The trick isn't trusting the AI more — it's trusting it less, in exactly the right place, so you can "
         "safely let it do far more everywhere else. That's the whole idea, and it's the step I think the one-shot "
         "skips. And to be clear — the gate covers what you CAN make deterministic: build, tests, scope, wiring. "
         "The genuinely subjective calls still go to the human review step. The gate doesn't replace judgment; it "
         "stops the agent from faking the part that can be checked. Put another way: the human review step is "
         "where you decide what 'done' means for each task; the deterministic gate just enforces that decision so "
         "the agent can't quietly redefine it mid-run.")
footer(s, 10)

# =====================================================================
# SLIDE 11 — See it for yourself (demo)
# =====================================================================
s = slide()
header(s, "See It for Yourself", label="LIVE DEMO")
tbox(s, 0.49, 1.0, 9.03, 0.62,
     [{"line": 1.08, "runs": [R("A tiny public demo — ", 13, SLATE, False), R("already broken down (6 tasks)", 13, BLUE, True),
                R(", ready to run.  Watch ", 13, SLATE, False), R("two tasks run in parallel", 13, ORANGE_DK, True),
                R(" and an ", 13, SLATE, False),
                R("AI resolve a real merge conflict behind a deterministic gate.", 13, SLATE, False)]}])
codebox(s, 0.49, 1.74, 9.03, 2.35,
    [[R("# 1. install the harness (latest preview)", 9.5, CODEMUT, False, MONO)],
     [R("dotnet tool install -g ServantSoftware.Guardrails ", 9.5, CODEFG, False, MONO),
      R("--version 1.0.0-preview.34", 9.5, CODEORG, False, MONO)],
     [R("", 6, CODEFG, False, MONO)],
     [R("# 2. clone the pre-broken-down demo", 9.5, CODEMUT, False, MONO)],
     [R("git clone https://github.com/DaveRMaltby/guardrails-texttools-demo", 9.5, CODEFG, False, MONO)],
     [R("cd guardrails-texttools-demo", 9.5, CODEFG, False, MONO)],
     [R("", 6, CODEFG, False, MONO)],
     [R("# 3. see the breakdown, preview the run, then run it", 9.5, CODEMUT, False, MONO)],
     [R("guardrails graph texttools", 9.5, CODEFG, False, MONO),
      R("            # writes diagram.html — the visible DAG", 9, CODEMUT, False, MONO)],
     [R("guardrails run   texttools ", 9.5, CODEFG, False, MONO),
      R("--dry-run", 9.5, CODEGRN, False, MONO),
      R("    # preview waves, zero cost", 9, CODEMUT, False, MONO)],
     [R("guardrails run   texttools", 9.5, CODEFG, False, MONO),
      R("              # execute — each task proves it worked", 9, CODEMUT, False, MONO)]])
foot = rect(s, 0.49, 4.26, 9.03, 0.72, fill=BLUE_TINT, line=BLUE_TINT2, line_w=1, rounded=True, radius=0.08)
settext(foot, [{"line": 1.05, "runs": [R("Needs .NET 8+, git, and the ", 11, SLATE, False),
               R("claude", 11, BLUE, True, MONO), R(" CLI.  Open ", 11, SLATE, False),
               R("texttools/diagram.html", 11, BLUE, True, MONO),
               R(" to see the breakdown you’d normally never get to review.", 11, SLATE, False)]}],
        anchor=MSO_ANCHOR.MIDDLE, ml=0.22, mr=0.2)
notes(s, "And if you want to feel it rather than take my word, here's a tiny public repo you can run yourself — "
         "I'll be running this exact one in the demo. It's already broken down, so you can skip straight to looking "
         "at the structure and running it. Install the tool, clone the repo, and the two commands that matter are "
         "'graph' — which opens that visible DAG, the thing the one-shot never shows you — and 'run', which "
         "executes it with every task proving itself. The dry-run flag previews the whole thing at zero cost. The "
         "fun part to watch: two tasks run in parallel, both edit the same file, and an AI resolves the merge "
         "conflict — but only after a deterministic re-verify says the merged result actually builds and passes. "
         "Try it on your own machine while I run it up here.")
footer(s, 11)

# =====================================================================
# SLIDE 12 — Back to the question
# =====================================================================
s = slide()
rect(s, 0, 0, 10, 5.625, fill=WHITE)
rect(s, 0, 0, 10, 0.16, fill=BLUE)
rect(s, 0, 0.16, 10, 0.04, fill=ORANGE)
tbox(s, 0.58, 1.15, 9.0, 0.95, [{"line": 1.02, "runs": [R("So — are we missing a step?", 38, BLUE, True)]}])
rect(s, 0.62, 2.12, 1.5, 0.05, fill=ORANGE)
tbox(s, 0.6, 2.34, 8.9, 1.4,
     [{"line": 1.18, "runs": [R("I think we might be.  Between ", 16, SLATE, False),
        R("“write a plan”", 16, DARK, True, ARIAL, True),
        R(" and ", 16, SLATE, False), R("“ship the diff”", 16, DARK, True, ARIAL, True),
        R(" there’s a step where a human should still see the breakdown, read the prompts, and trust a ", 16, SLATE, False),
        R("gate", 16, BLUE, True), R(", not a ", 16, SLATE, False), R("claim.", 16, ORANGE_DK, True)]}])
np = rect(s, 0.6, 3.95, 8.85, 0.62, fill=ORG_TINT, line=RGBColor(0xF6,0xD8,0xBE), line_w=1, rounded=True, radius=0.1)
settext(np, [{"align": PP_ALIGN.CENTER, "runs": [R("I’m not prescribing.  Just putting it on the table.", 15, ORANGE_DK, True)]}],
        anchor=MSO_ANCHOR.MIDDLE)
tbox(s, 0.6, 4.78, 9.0, 0.5,
     [{"runs": [R("Demo:  ", 10.5, MUTED, True), R("github.com/DaveRMaltby/guardrails-texttools-demo", 10.5, SLATE, False, MONO),
                R("      Project:  ", 10.5, MUTED, True), R("github.com/Servant-Software-LLC/Guardrails", 10.5, SLATE, False, MONO)]}])
notes(s, "So I'll end where I started, with the question. Are we missing a step in agentic engineering? I think we "
         "might be — a step between writing the plan and shipping the diff, where a human still gets to see the "
         "breakdown, read the prompts, and rely on a deterministic gate instead of the agent's word. I built one "
         "answer to that, but I'm genuinely not here to prescribe it. I just wanted to put the question on the "
         "table, because I don't think we talk about that middle step enough. Happy to run the demo, and happier to "
         "argue about whether the step is real. Thank you.")
footer(s, 12)

import os
out = os.environ.get("DECK_OUT", r"C:\Dev AI\Guardrails\docs\presentations\tosca-learns-2026-06-30.pptx")
prs.save(out)
print("Saved:", out, "with", len(prs.slides._sldIdLst), "slides")
